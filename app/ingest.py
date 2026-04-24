import argparse
import json
from pathlib import Path
import sys

if __package__ is None or __package__ == "":
    sys.path.append(str(Path(__file__).resolve().parent.parent))

from lightrag.base import DocStatus

from app.rag import get_rag


async def _ensure_rag_initialized(rag) -> None:
    await rag.initialize_storages()


def extract_text_from_pdf(filepath: str) -> str:
    try:
        import PyPDF2
    except ImportError as exc:
        raise RuntimeError("PyPDF2 is required to ingest PDF files.") from exc

    text = ""
    with open(filepath, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text += page_text + "\n"
    return text


def extract_text_from_docx(filepath: str) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to ingest DOCX files.") from exc

    doc = Document(filepath)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def extract_text_from_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to ingest PPTX files.") from exc

    prs = Presentation(filepath)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
    return "\n".join(lines)


def extract_text_from_xlsx(filepath: str) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to ingest XLSX/XLS files.") from exc

    wb = openpyxl.load_workbook(filepath, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def extract_text_from_csv(filepath: str) -> str:
    import csv

    lines: list[str] = []
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = "\t".join(row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def extract_text(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    if ext == ".docx":
        return extract_text_from_docx(filepath)
    if ext in (".txt", ".md"):
        return extract_text_from_txt(filepath)
    if ext == ".pptx":
        return extract_text_from_pptx(filepath)
    if ext in (".xlsx", ".xls"):
        return extract_text_from_xlsx(filepath)
    if ext == ".csv":
        return extract_text_from_csv(filepath)
    raise ValueError(f"Unsupported file type: {ext}")


def _normalize_path(path_value: str) -> str:
    return str(Path(path_value).expanduser().resolve())


def _matches_file(existing_file_path: str, new_file_path: str, match_mode: str) -> bool:
    if not existing_file_path:
        return False

    existing_norm = _normalize_path(existing_file_path)
    new_norm = _normalize_path(new_file_path)

    if match_mode == "exact":
        return existing_norm == new_norm
    if match_mode == "basename":
        return Path(existing_norm).name.lower() == Path(new_norm).name.lower()
    raise ValueError(f"Unsupported match_mode: {match_mode}")


async def _all_docs(rag) -> dict:
    docs: dict = {}
    for status in DocStatus:
        docs.update(await rag.get_docs_by_status(status))
    return docs


async def list_docs() -> list[dict]:
    rag = get_rag()
    await _ensure_rag_initialized(rag)
    docs = await _all_docs(rag)
    items: list[dict] = []
    for doc_id, status in docs.items():
        items.append(
            {
                "doc_id": doc_id,
                "status": str(getattr(status, "status", "")),
                "file_path": getattr(status, "file_path", ""),
                "updated_at": getattr(status, "updated_at", ""),
            }
        )
    return sorted(items, key=lambda item: item["updated_at"], reverse=True)


async def delete_doc_id(doc_id: str) -> dict:
    rag = get_rag()
    await _ensure_rag_initialized(rag)
    result = await rag.adelete_by_doc_id(doc_id, delete_llm_cache=True)
    await rag.aclear_cache()
    return {
        "doc_id": doc_id,
        "status": result.status,
        "message": result.message,
        "status_code": result.status_code,
        "file_path": result.file_path,
    }


async def delete_existing_versions(
    filepath: str,
    match_mode: str = "basename",
    delete_llm_cache: bool = True,
) -> list[dict]:
    rag = get_rag()
    await _ensure_rag_initialized(rag)
    docs = await _all_docs(rag)
    deleted: list[dict] = []

    for doc_id, status in docs.items():
        doc_file_path = getattr(status, "file_path", "")
        if not _matches_file(doc_file_path, filepath, match_mode):
            continue

        result = await rag.adelete_by_doc_id(
            doc_id,
            delete_llm_cache=delete_llm_cache,
        )
        deleted.append(
            {
                "doc_id": doc_id,
                "file_path": doc_file_path,
                "status": result.status,
                "message": result.message,
                "status_code": result.status_code,
            }
        )

    return deleted


async def ingest_file(
    filepath: str,
    replace_existing: bool = False,
    match_mode: str = "basename",
) -> dict:
    rag = get_rag()
    await _ensure_rag_initialized(rag)
    filename = Path(filepath).name

    try:
        text = extract_text(filepath)
        if not text.strip():
            return {
                "success": False,
                "filename": filename,
                "error": "File appears to be empty or unreadable",
            }

        deleted_docs: list[dict] = []
        if replace_existing:
            deleted_docs = await delete_existing_versions(
                filepath=filepath,
                match_mode=match_mode,
                delete_llm_cache=True,
            )
            if deleted_docs:
                await rag.aclear_cache()

        track_id = await rag.ainsert(text, file_paths=[filepath])
        return {
            "success": True,
            "filename": filename,
            "chars": len(text),
            "track_id": track_id,
            "replaced_docs": deleted_docs,
        }
    except Exception as exc:
        return {"success": False, "filename": filename, "error": str(exc)}


async def replace_file(filepath: str, match_mode: str = "basename") -> dict:
    return await ingest_file(
        filepath=filepath,
        replace_existing=True,
        match_mode=match_mode,
    )


async def ingest_directory(
    directory: str,
    replace_existing: bool = False,
    match_mode: str = "basename",
) -> list:
    dir_path = Path(directory)
    if not dir_path.exists():
        return []

    results = []
    supported = [".pdf", ".docx", ".txt", ".md", ".pptx", ".xlsx", ".xls", ".csv"]
    for f in dir_path.iterdir():
        if f.suffix.lower() not in supported:
            continue
        result = await ingest_file(
            str(f),
            replace_existing=replace_existing,
            match_mode=match_mode,
        )
        results.append(result)
    return results


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Ingest files into LightRAG storage.")
    sub = parser.add_subparsers(dest="command", required=True)

    ingest_file_parser = sub.add_parser("ingest-file", help="Ingest one file.")
    ingest_file_parser.add_argument("filepath")
    ingest_file_parser.add_argument(
        "--replace-existing",
        action="store_true",
        help="Delete older versions of the same file first.",
    )
    ingest_file_parser.add_argument(
        "--match-mode",
        choices=["basename", "exact"],
        default="basename",
        help="How to detect older versions.",
    )

    replace_file_parser = sub.add_parser(
        "replace-file",
        help="Delete old versions and ingest this file.",
    )
    replace_file_parser.add_argument("filepath")
    replace_file_parser.add_argument(
        "--match-mode",
        choices=["basename", "exact"],
        default="basename",
        help="How to detect older versions.",
    )

    cleanup_file_parser = sub.add_parser(
        "cleanup-file",
        help="Delete existing versions of a file without ingesting a new one.",
    )
    cleanup_file_parser.add_argument("filepath")
    cleanup_file_parser.add_argument(
        "--match-mode",
        choices=["basename", "exact"],
        default="basename",
        help="How to detect older versions.",
    )

    delete_doc_id_parser = sub.add_parser(
        "delete-doc-id",
        help="Delete one indexed document by its doc_id.",
    )
    delete_doc_id_parser.add_argument("doc_id")

    sub.add_parser("list-docs", help="List indexed documents in rag storage.")

    ingest_dir_parser = sub.add_parser("ingest-dir", help="Ingest all supported files in a directory.")
    ingest_dir_parser.add_argument("directory")
    ingest_dir_parser.add_argument(
        "--replace-existing",
        action="store_true",
        help="Delete older versions of each matching file first.",
    )
    ingest_dir_parser.add_argument(
        "--match-mode",
        choices=["basename", "exact"],
        default="basename",
        help="How to detect older versions.",
    )

    return parser


if __name__ == "__main__":
    import asyncio

    args = _build_parser().parse_args()

    if args.command == "ingest-file":
        output = asyncio.run(
            ingest_file(
                filepath=args.filepath,
                replace_existing=args.replace_existing,
                match_mode=args.match_mode,
            )
        )
    elif args.command == "replace-file":
        output = asyncio.run(
            replace_file(
                filepath=args.filepath,
                match_mode=args.match_mode,
            )
        )
    elif args.command == "ingest-dir":
        output = asyncio.run(
            ingest_directory(
                directory=args.directory,
                replace_existing=args.replace_existing,
                match_mode=args.match_mode,
            )
        )
    elif args.command == "cleanup-file":
        output = asyncio.run(
            delete_existing_versions(
                filepath=args.filepath,
                match_mode=args.match_mode,
                delete_llm_cache=True,
            )
        )
    elif args.command == "list-docs":
        output = asyncio.run(list_docs())
    elif args.command == "delete-doc-id":
        output = asyncio.run(delete_doc_id(args.doc_id))
    else:
        raise ValueError(f"Unsupported command: {args.command}")

    print(json.dumps(output, indent=2, ensure_ascii=False))
